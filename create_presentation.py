#!/usr/bin/env python3
"""
Generates a PowerPoint presentation for the Real-Time Release Management System project.
Run: pip install python-pptx && python create_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_content_slide(prs, title, content_points):
    """Add a content slide with bullet points."""
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, point in enumerate(content_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
    return slide

def add_two_column_slide(prs, title, left_content, right_content):
    """Add a slide with two columns."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(5.5))
    left_tf = left_box.text_frame
    left_tf.word_wrap = True
    for i, point in enumerate(left_content):
        if i == 0:
            p = left_tf.paragraphs[0]
        else:
            p = left_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(14)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(5.5))
    right_tf = right_box.text_frame
    right_tf.word_wrap = True
    for i, point in enumerate(right_content):
        if i == 0:
            p = right_tf.paragraphs[0]
        else:
            p = right_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(14)
    
    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs, 
        "Real-Time Release Management System",
        "CS544 Enterprise Architecture - Project Presentation"
    )

    # Slide 2: Project Overview
    add_content_slide(prs, "Project Overview", [
        "A comprehensive platform for managing software releases with strict workflow enforcement",
        "Designed for fast-growing software companies",
        "Key capabilities:",
        "  • Full release lifecycle: task assignment → completion → hotfixes",
        "  • Sequential task execution and single task-in-process rule per developer",
        "  • Event-driven architecture using Apache Kafka",
        "  • Real-time collaboration via Server-Sent Events (SSE)",
        "  • AI-powered developer assistance (Ollama integration)",
        "  • Email notifications when tasks are assigned",
        "  • Observability with Prometheus and Grafana"
    ])

    # Slide 3: System Architecture
    add_content_slide(prs, "System Architecture", [
        "Microservices Architecture:",
        "",
        "1. Release Service (Port 8080)",
        "   • Core business logic, workflow management, task tracking",
        "   • Forum/discussion threads on tasks",
        "   • AI chat with Ollama for developer assistance",
        "   • Tech: Spring Boot, MongoDB, Kafka, WebFlux (SSE)",
        "",
        "2. Notification Service (Port 8081)",
        "   • Consumes Kafka events for task assignments, hotfixes, stale tasks",
        "   • Sends real emails via Gmail SMTP",
        "   • Maintains notification audit logs in MongoDB",
        "   • Tech: Spring Boot, Kafka, MongoDB, JavaMail"
    ])

    # Slide 4: Frontend & Infrastructure
    add_content_slide(prs, "Frontend & Infrastructure", [
        "3. Frontend (Port 5173)",
        "   • React + TypeScript + Vite",
        "   • Pages: Login, Register, Dashboard, Releases (Admin), My Tasks (Developer)",
        "   • Activity Feed (real-time SSE), AI Chat",
        "   • JWT authentication",
        "",
        "Infrastructure (Docker Compose):",
        "   • Kafka + Zookeeper (message broker)",
        "   • MongoDB (release_db, notification_db)",
        "   • Ollama (local LLM for AI chat)",
        "   • Prometheus + Grafana (monitoring)"
    ])

    # Slide 5: Workflow Enforcement
    add_content_slide(prs, "Workflow Enforcement Rules", [
        "1. Single In-Process Rule",
        "   • Only ONE task per developer can be IN_PROCESS at any time",
        "   • Prevents developers from juggling multiple active tasks",
        "   • Enforced in ReleaseService.startTask()",
        "",
        "2. Sequential Task Execution",
        "   • Tasks must be completed in release order",
        "   • Cannot start task N until task N-1 is COMPLETED",
        "   • Enforced before allowing task start",
        "",
        "3. Hotfix Logic",
        "   • Adding a task to a COMPLETED release automatically re-opens it",
        "   • Publishes HotfixTaskAddedEvent to Kafka",
        "   • Developer receives urgent email notification"
    ])

    # Slide 6: Event-Driven Architecture (Kafka)
    add_content_slide(prs, "Event-Driven Architecture with Kafka", [
        "Topic: task-events",
        "   • assigned – when a task is added and assigned to a developer",
        "   • hotfix – when a task is added to a completed release",
        "   • stale – when a task has been IN_PROGRESS for 48+ hours (scheduler)",
        "   • completed – when a developer completes a task",
        "",
        "Topic: system-events",
        "   • error – system health monitor publishes errors",
        "",
        "Flow: Release Service (producer) → Kafka → Notification Service (consumer)",
        "Notification Service sends emails and logs to MongoDB"
    ])

    # Slide 7: Authentication & User Management
    add_content_slide(prs, "Authentication & User Management", [
        "JWT-based authentication",
        "   • POST /auth/register – Register with username, email, password, role",
        "   • POST /auth/login – Returns JWT token",
        "   • GET /auth/me – Current user info",
        "",
        "Roles: ADMIN, DEVELOPER",
        "   • ADMIN: Create releases, add tasks, assign developers, complete releases",
        "   • DEVELOPER: View assigned tasks, start/complete tasks, add comments, use AI chat",
        "",
        "Email field: Stored per user for task-assignment notifications",
        "Password: bcrypt encrypted"
    ])

    # Slide 8: Real-Time Features
    add_content_slide(prs, "Real-Time Features", [
        "1. Activity Feed (SSE)",
        "   • GET /activity/stream – Server-Sent Events",
        "   • Live updates: task started, completed, hotfix added, comments, etc.",
        "   • Frontend subscribes and displays in real time",
        "",
        "2. AI Chat Assistant",
        "   • POST /chat/session – Create session",
        "   • POST /chat/{sessionId}/message – Send message",
        "   • Backend calls Ollama (local LLM) for responses",
        "   • Context-aware developer assistance",
        "",
        "3. Forum/Discussion",
        "   • Nested comments on tasks (Reddit-style threading)",
        "   • Add comment, reply to comment"
    ])

    # Slide 9: Notification & Email Flow
    add_content_slide(prs, "Notification & Email Flow", [
        "When a task is assigned:",
        "   1. Admin adds task via Releases page, selects developer (username)",
        "   2. ReleaseService looks up developer email from UserRepository",
        "   3. Publishes TaskAssignedEvent (taskId, developerId, developerEmail) to Kafka",
        "   4. Notification Service consumes event",
        "   5. Sends email to developer's registered email via Gmail SMTP",
        "   6. Logs to NotificationLog in MongoDB",
        "",
        "Same flow for: HotfixTaskAddedEvent, StaleTaskDetectedEvent",
        "Stale Task Scheduler runs hourly, detects tasks IN_PROGRESS for 48+ hours"
    ])

    # Slide 10: API Endpoints Summary
    add_content_slide(prs, "API Endpoints Summary", [
        "Auth: /auth/register, /auth/login, /auth/me",
        "Releases: POST /releases, GET /releases, POST /releases/{id}/tasks, PATCH /releases/{id}/complete",
        "Tasks: GET /tasks/my, PATCH /tasks/{id}/start, PATCH /tasks/{id}/complete",
        "Comments: POST /tasks/{id}/comments, GET /tasks/{id}/comments, POST /comments/{id}/reply",
        "Activity: GET /activity/stream (SSE)",
        "Chat: POST /chat/session, POST /chat/{sessionId}/message, GET /chat/{sessionId}/history",
        "Test: GET /test/email?to=... (notification service)"
    ])

    # Slide 11: Technology Stack
    add_content_slide(prs, "Technology Stack", [
        "Backend: Java 21, Spring Boot 3.2, Spring Data MongoDB, Spring Kafka, Spring Security (JWT)",
        "Frontend: React 18, TypeScript, Vite",
        "Database: MongoDB (document store)",
        "Message Broker: Apache Kafka",
        "AI: Ollama (local LLM)",
        "Email: Gmail SMTP (Spring Mail)",
        "Monitoring: Prometheus, Grafana",
        "Containerization: Docker, Docker Compose"
    ])

    # Slide 12: How to Run
    add_content_slide(prs, "How to Run the Project", [
        "1. Prerequisites: Docker, Docker Compose, Maven",
        "2. Build: mvn clean package -DskipTests",
        "3. Run: docker-compose up -d",
        "4. Access:",
        "   • Frontend: http://localhost:5173",
        "   • Release API: http://localhost:8080",
        "   • Notification API: http://localhost:8081",
        "   • Grafana: http://localhost:3000",
        "   • Prometheus: http://localhost:9090",
        "5. Gmail credentials in notification-service/application.yml for email delivery"
    ])

    # Slide 13: Thank You
    add_title_slide(prs, "Thank You", "Questions?")

    # Save
    output_path = "ReleaseManagementSystem_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
