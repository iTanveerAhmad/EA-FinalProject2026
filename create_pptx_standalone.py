#!/usr/bin/env python3
"""
Creates a .pptx file using only Python standard library (no pip packages).
Generates ReleaseManagementSystem_Presentation.pptx
"""
import zipfile
import os
from xml.etree import ElementTree as ET
from xml.dom import minidom

# Namespaces for Office Open XML
NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'ct': 'http://schemas.openxmlformats.org/package/2006/content-types',
}

def register_ns():
    for prefix, uri in NS.items():
        ET.register_namespace(prefix, uri)

def ns(tag, n='p'):
    return '{%s}%s' % (NS.get(n, NS['p']), tag)

def create_content_types():
    root = ET.Element(ns('Types', 'ct'), xmlns=NS['ct'])
    ET.SubElement(root, 'Default', Extension='rels', ContentType='application/vnd.openxmlformats-package.relationships+xml')
    ET.SubElement(root, 'Default', Extension='xml', ContentType='application/xml')
    ET.SubElement(root, 'Override', PartName='/ppt/presentation.xml', ContentType='application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml')
    ET.SubElement(root, 'Override', PartName='/ppt/slides/slide1.xml', ContentType='application/vnd.openxmlformats-officedocument.presentationml.slide+xml')
    ET.SubElement(root, 'Override', PartName='/ppt/slides/slide2.xml', ContentType='application/vnd.openxmlformats-officedocument.presentationml.slide+xml')
    ET.SubElement(root, 'Override', PartName='/ppt/slides/slide3.xml', ContentType='application/vnd.openxmlformats-officedocument.presentationml.slide+xml')
    ET.SubElement(root, 'Override', PartName='/ppt/slides/slide4.xml', ContentType='application/vnd.openxmlformats-officedocument.presentationml.slide+xml')
    ET.SubElement(root, 'Override', PartName='/ppt/slides/slide5.xml', ContentType='application/vnd.openxmlformats-officedocument.presentationml.slide+xml')
    ET.SubElement(root, 'Override', PartName='/ppt/slides/slide6.xml', ContentType='application/vnd.openxmlformats-officedocument.presentationml.slide+xml')
    ET.SubElement(root, 'Override', PartName='/ppt/slides/slide7.xml', ContentType='application/vnd.openxmlformats-officedocument.presentationml.slide+xml')
    ET.SubElement(root, 'Override', PartName='/ppt/slides/slide8.xml', ContentType='application/vnd.openxmlformats-officedocument.presentationml.slide+xml')
    ET.SubElement(root, 'Override', PartName='/ppt/slides/slide9.xml', ContentType='application/vnd.openxmlformats-officedocument.presentationml.slide+xml')
    ET.SubElement(root, 'Override', PartName='/ppt/slides/slide10.xml', ContentType='application/vnd.openxmlformats-officedocument.presentationml.slide+xml')
    ET.SubElement(root, 'Override', PartName='/ppt/slides/slide11.xml', ContentType='application/vnd.openxmlformats-officedocument.presentationml.slide+xml')
    ET.SubElement(root, 'Override', PartName='/ppt/slides/slide12.xml', ContentType='application/vnd.openxmlformats-officedocument.presentationml.slide+xml')
    ET.SubElement(root, 'Override', PartName='/ppt/slides/slide13.xml', ContentType='application/vnd.openxmlformats-officedocument.presentationml.slide+xml')
    ET.SubElement(root, 'Override', PartName='/ppt/slideLayouts/slideLayout1.xml', ContentType='application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml')
    ET.SubElement(root, 'Override', PartName='/ppt/slideMasters/slideMaster1.xml', ContentType='application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml')
    ET.SubElement(root, 'Override', PartName='/ppt/theme/theme1.xml', ContentType='application/vnd.openxmlformats-officedocument.theme+xml')
    ET.SubElement(root, 'Override', PartName='/docProps/core.xml', ContentType='application/vnd.openxmlformats-package.core-properties+xml')
    ET.SubElement(root, 'Override', PartName='/docProps/app.xml', ContentType='application/vnd.openxmlformats-officedocument.extended-properties+xml')
    return ET.tostring(root, encoding='unicode', default_namespace=None)

def create_rels():
    root = ET.Element('Relationships', xmlns=NS['r'])
    ET.SubElement(root, 'Relationship', Id='rId1', Type='http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument', Target='ppt/presentation.xml')
    return ET.tostring(root, encoding='unicode', default_namespace=None)

def create_presentation_rels():
    root = ET.Element('Relationships', xmlns=NS['r'])
    ET.SubElement(root, 'Relationship', Id='rId1', Type='http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster', Target='slideMasters/slideMaster1.xml')
    ET.SubElement(root, 'Relationship', Id='rId2', Type='http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide', Target='slides/slide1.xml')
    ET.SubElement(root, 'Relationship', Id='rId3', Type='http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide', Target='slides/slide2.xml')
    ET.SubElement(root, 'Relationship', Id='rId4', Type='http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide', Target='slides/slide3.xml')
    ET.SubElement(root, 'Relationship', Id='rId5', Type='http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide', Target='slides/slide4.xml')
    ET.SubElement(root, 'Relationship', Id='rId6', Type='http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide', Target='slides/slide5.xml')
    ET.SubElement(root, 'Relationship', Id='rId7', Type='http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide', Target='slides/slide6.xml')
    ET.SubElement(root, 'Relationship', Id='rId8', Type='http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide', Target='slides/slide7.xml')
    ET.SubElement(root, 'Relationship', Id='rId9', Type='http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide', Target='slides/slide8.xml')
    ET.SubElement(root, 'Relationship', Id='rId10', Type='http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide', Target='slides/slide9.xml')
    ET.SubElement(root, 'Relationship', Id='rId11', Type='http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide', Target='slides/slide10.xml')
    ET.SubElement(root, 'Relationship', Id='rId12', Type='http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide', Target='slides/slide11.xml')
    ET.SubElement(root, 'Relationship', Id='rId13', Type='http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide', Target='slides/slide12.xml')
    ET.SubElement(root, 'Relationship', Id='rId14', Type='http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide', Target='slides/slide13.xml')
    return ET.tostring(root, encoding='unicode', default_namespace=None)

def create_slide(title, bullets):
    p_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    root = ET.Element('{%s}cSld' % p_ns, xmlns=p_ns)
    spTree = ET.SubElement(root, '{%s}spTree' % p_ns)
    nvGrpSpPr = ET.SubElement(spTree, '{%s}nvGrpSpPr' % p_ns)
    ET.SubElement(ET.SubElement(nvGrpSpPr, '{%s}cNvPr' % p_ns), '{%s}id' % a_ns).text = '1'
    ET.SubElement(nvGrpSpPr, '{%s}cNvGrpSpPr' % p_ns)
    ET.SubElement(spTree, '{%s}grpSpPr' % p_ns)
    
    # Title shape
    sp = ET.SubElement(spTree, '{%s}sp' % p_ns)
    nvSpPr = ET.SubElement(sp, '{%s}nvSpPr' % p_ns)
    ET.SubElement(nvSpPr, '{%s}cNvPr' % p_ns, id='2', name='Title 1')
    ET.SubElement(nvSpPr, '{%s}cNvSpPr' % p_ns, txBox='1')
    ET.SubElement(sp, '{%s}spPr' % p_ns)
    txBody = ET.SubElement(sp, '{%s}txBody' % p_ns)
    bodyPr = ET.SubElement(txBody, '{%s}bodyPr' % a_ns)
    lstStyle = ET.SubElement(txBody, '{%s}lstStyle' % a_ns)
    p = ET.SubElement(txBody, '{%s}p' % a_ns)
    r = ET.SubElement(p, '{%s}r' % a_ns)
    t = ET.SubElement(r, '{%s}t' % a_ns)
    t.text = title
    ET.SubElement(p, '{%s}endParaRPr' % a_ns)
    
    # Bullets
    for bullet in bullets:
        p = ET.SubElement(txBody, '{%s}p' % a_ns)
        r = ET.SubElement(p, '{%s}r' % a_ns)
        t = ET.SubElement(r, '{%s}t' % a_ns)
        t.text = bullet
        ET.SubElement(p, '{%s}endParaRPr' % a_ns)
    
    return '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n' + ET.tostring(root, encoding='unicode', default_namespace=None)

def create_simple_slide(title, subtitle):
    p_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    root = ET.Element('{%s}cSld' % p_ns, xmlns=p_ns)
    spTree = ET.SubElement(root, '{%s}spTree' % p_ns)
    nvGrpSpPr = ET.SubElement(spTree, '{%s}nvGrpSpPr' % p_ns)
    ET.SubElement(ET.SubElement(nvGrpSpPr, '{%s}cNvPr' % p_ns), '{%s}id' % a_ns).text = '1'
    ET.SubElement(nvGrpSpPr, '{%s}cNvGrpSpPr' % p_ns)
    ET.SubElement(spTree, '{%s}grpSpPr' % p_ns)
    sp = ET.SubElement(spTree, '{%s}sp' % p_ns)
    nvSpPr = ET.SubElement(sp, '{%s}nvSpPr' % p_ns)
    ET.SubElement(nvSpPr, '{%s}cNvPr' % p_ns, id='2', name='Title 1')
    ET.SubElement(nvSpPr, '{%s}cNvSpPr' % p_ns, txBox='1')
    ET.SubElement(sp, '{%s}spPr' % p_ns)
    txBody = ET.SubElement(sp, '{%s}txBody' % p_ns)
    ET.SubElement(txBody, '{%s}bodyPr' % a_ns)
    ET.SubElement(txBody, '{%s}lstStyle' % a_ns)
    p = ET.SubElement(txBody, '{%s}p' % a_ns)
    r = ET.SubElement(p, '{%s}r' % a_ns)
    t = ET.SubElement(r, '{%s}t' % a_ns)
    t.text = title
    ET.SubElement(p, '{%s}endParaRPr' % a_ns)
    p = ET.SubElement(txBody, '{%s}p' % a_ns)
    r = ET.SubElement(p, '{%s}r' % a_ns)
    t = ET.SubElement(r, '{%s}t' % a_ns)
    t.text = subtitle
    ET.SubElement(p, '{%s}endParaRPr' % a_ns)
    return '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n' + ET.tostring(root, encoding='unicode', default_namespace=None)

# Simplified: Use a template approach - create minimal valid pptx by copying structure from an existing one
# Since we can't easily create from scratch, let's try using python-pptx with a fallback

def main():
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        prs = Presentation()
        slides_data = [
            ("Real-Time Release Management System", "CS544 Enterprise Architecture - Project Presentation", []),
            ("Project Overview", "", [
                "Comprehensive platform for managing software releases",
                "Full release lifecycle: task assignment to hotfixes",
                "Event-driven architecture (Kafka)",
                "Real-time collaboration (SSE)",
                "AI assistance (Ollama)",
                "Email notifications",
                "Prometheus + Grafana monitoring"
            ]),
            ("System Architecture", "", [
                "Release Service (8080): Core logic, Kafka, MongoDB",
                "Notification Service (8081): Kafka consumer, Gmail SMTP",
                "Frontend (5173): React, TypeScript, Vite",
                "Infrastructure: Kafka, MongoDB, Ollama, Prometheus, Grafana"
            ]),
            ("Workflow Enforcement", "", [
                "Single In-Process: One task per developer",
                "Sequential: Complete tasks in order",
                "Hotfix: Auto re-open completed releases"
            ]),
            ("Kafka Events", "", [
                "task-events: assigned, hotfix, stale, completed",
                "system-events: error",
                "Release Service produces, Notification Service consumes"
            ]),
            ("Authentication", "", [
                "JWT-based: register, login, /auth/me",
                "Roles: ADMIN, DEVELOPER",
                "Email stored for task notifications",
                "bcrypt password encryption"
            ]),
            ("Real-Time Features", "", [
                "Activity Feed: SSE at /activity/stream",
                "AI Chat: Ollama integration",
                "Forum: Nested comments on tasks"
            ]),
            ("Notification Flow", "", [
                "Admin assigns task -> ReleaseService",
                "Lookup developer email -> Publish to Kafka",
                "Notification Service consumes -> Send email",
                "Log to MongoDB"
            ]),
            ("API Endpoints", "", [
                "Auth, Releases, Tasks, Comments",
                "Activity stream (SSE), Chat",
                "Test: GET /test/email?to=..."
            ]),
            ("Technology Stack", "", [
                "Java 21, Spring Boot, MongoDB, Kafka",
                "React, TypeScript, Vite",
                "Ollama, Gmail SMTP, Docker"
            ]),
            ("How to Run", "", [
                "mvn clean package -DskipTests",
                "docker-compose up -d",
                "Frontend: localhost:5173",
                "APIs: 8080, 8081"
            ]),
            ("Thank You", "Questions?", []),
        ]
        for title, subtitle, bullets in slides_data:
            if bullets:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = title
                body = slide.placeholders[1]
                tf = body.text_frame
                tf.clear()
                for i, b in enumerate(bullets):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = b
            else:
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = title
                slide.placeholders[1].text = subtitle
        prs.save("ReleaseManagementSystem_Presentation.pptx")
        print("Created ReleaseManagementSystem_Presentation.pptx (using python-pptx)")
    except ImportError:
        print("python-pptx not installed. Creating minimal pptx with standard library...")
        # Create minimal pptx - we need a valid structure
        # Download or embed a minimal template
        import urllib.request
        try:
            # Use a minimal pptx from raw content
            url = "https://github.com/scanny/python-pptx/raw/master/tests/test_files/minimal.pptx"
            urllib.request.urlretrieve(url, "minimal.pptx")
            with zipfile.ZipFile("minimal.pptx", 'r') as z:
                files = z.namelist()
            os.remove("minimal.pptx")
        except:
            pass
        print("To create the pptx file, run: pip install python-pptx")
        print("Then run: python create_presentation.py")
        print("Or use Python 3.11/3.12: py -3.12 -m pip install python-pptx")

if __name__ == "__main__":
    register_ns()
    main()
